import os
import io
import zipfile
import pytesseract
from PIL import Image
import pdfplumber
import docx
import openpyxl
from pptx import Presentation
import pandas as pd

def extract_text_from_file(file_obj, filename):
    """
    Extracts text from various file formats, including zip archives.
    """
    ext = os.path.splitext(filename)[1].lower()
    text = ""
    try:
        if ext == '.zip':
            with zipfile.ZipFile(file_obj, 'r') as z:
                for inner_filename in z.namelist():
                    # Skip directories and hidden files like __MACOSX
                    if not inner_filename.endswith('/') and not inner_filename.startswith('__MACOSX/'):
                        with z.open(inner_filename) as inner_file:
                            inner_file_bytes = io.BytesIO(inner_file.read())
                            extracted = extract_text_from_file(inner_file_bytes, inner_filename)
                            if extracted and not extracted.startswith("Unsupported file type") and not extracted.startswith("Error processing"):
                                text += f"\n--- Arquivo no pacote: {inner_filename} ---\n"
                                text += extracted + "\n"
            return text.strip()
            
        elif ext == '.pdf':
            with pdfplumber.open(file_obj) as pdf:
                for page in pdf.pages:
                    # Extract regular text
                    extracted = page.extract_text()
                    if extracted:
                        text += extracted + "\n"
                    else:
                        # Fallback to OCR if PDF page is an image
                        img = page.to_image().original
                        text += pytesseract.image_to_string(img, lang='por+eng') + "\n"

                    # Extract tables as markdown
                    tables = page.extract_tables()
                    for table in tables:
                        if table:
                            df = pd.DataFrame(table[1:], columns=table[0])
                            # Clean up None/NaN values
                            df = df.fillna("")
                            text += "\n" + df.to_markdown(index=False) + "\n"

        elif ext in ['.docx', '.doc']:
            doc = docx.Document(file_obj)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif ext in ['.xlsx', '.xls']:
            xls = pd.ExcelFile(file_obj)
            for sheet_name in xls.sheet_names:
                text += f"\n### Sheet: {sheet_name}\n"
                df = pd.read_excel(xls, sheet_name=sheet_name)
                df = df.fillna("")
                text += df.to_markdown(index=False) + "\n"
        elif ext in ['.pptx', '.ppt']:
            prs = Presentation(file_obj)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif ext in ['.png', '.jpg', '.jpeg', '.webp']:
            img = Image.open(file_obj)
            text += pytesseract.image_to_string(img, lang='por+eng') + "\n"
        elif ext in ['.txt', '.csv']:
            # For txt or csv we assume it's read as string or bytes
            if isinstance(file_obj, io.BytesIO):
                text += file_obj.read().decode('utf-8', errors='ignore')
            elif isinstance(file_obj, str):
                with open(file_obj, 'r', encoding='utf-8', errors='ignore') as f:
                    text += f.read()
            else:
                 text += file_obj.read().decode('utf-8', errors='ignore')
        else:
            text = f"Unsupported file type: {ext}"
    except Exception as e:
        text = f"Error processing {filename}: {str(e)}"
    
    return text.strip()
